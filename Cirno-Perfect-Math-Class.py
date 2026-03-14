import sys
import os
import json
from pptx import Presentation
from openai import OpenAI
from PyQt6.QtWidgets import (QApplication, QMainWindow, QWidget, QVBoxLayout,
                             QHBoxLayout, QPushButton, QTextEdit, QFileDialog,
                             QLineEdit, QLabel, QMessageBox, QFrame, QComboBox)
from PyQt6.QtWebEngineWidgets import QWebEngineView
from PyQt6.QtCore import QThread, pyqtSignal

# ================= 核心：防闪退 & 直连机制 =================
os.environ["QTWEBENGINE_DISABLE_GPU"] = "1"
os.environ["QTWEBENGINE_CHROMIUM_FLAGS"] = "--disable-gpu --no-sandbox"

for proxy_env in ['http_proxy', 'https_proxy', 'HTTP_PROXY', 'HTTPS_PROXY', 'all_proxy', 'ALL_PROXY']:
    if proxy_env in os.environ:
        del os.environ[proxy_env]

def get_resource_path(relative_path):
    if getattr(sys, 'frozen', False):
        base_path = os.path.dirname(sys.executable)
    else:
        base_path = os.path.dirname(os.path.abspath(__file__))
    return os.path.join(base_path, relative_path).replace("\\", "/")


# ❄️ 多线程后台大脑 (负责不卡死界面)
class AIWorker(QThread):
    finished_signal = pyqtSignal(str)  
    error_signal = pyqtSignal(str)     

    def __init__(self, api_key, base_url, model_name, messages):
        super().__init__()
        self.api_key = api_key
        self.base_url = base_url
        self.model_name = model_name
        self.messages = messages

    def run(self):
        try:
            client = OpenAI(api_key=self.api_key, base_url=self.base_url)
            response = client.chat.completions.create(
                model=self.model_name,
                messages=self.messages
            )
            reply = response.choices[0].message.content
            self.finished_signal.emit(reply)
        except Exception as e:
            self.error_signal.emit(str(e))


class PPTAIAssistant(QMainWindow):
    def __init__(self):
        super().__init__()
        # ❄️ 纯正二次元标题
        self.setWindowTitle("❄️ 《琪露诺的完美算数教室 v1.0》 - 最强天才⑨的教研室 ❄️")
        self.resize(1280, 900)
        
        self.ppt_content = ""
        self.chat_history = []
        self.api_key = ""
        self.zoom_factor = 1.0

        self.apply_stylesheet()
        self.init_ui()

    def apply_stylesheet(self):
        bg_path = get_resource_path("cirno.jpg")
        self.setStyleSheet(f"""
            QMainWindow {{ background-color: #E0F2FE; border-image: url('{bg_path}'); }}
            .PanelFrame {{ background-color: rgba(255, 255, 255, 0.9); border-radius: 10px; border: 2px solid #BAE6FD; }}
            
            QLineEdit, QTextEdit, QComboBox {{ 
                background-color: rgba(240, 249, 255, 0.8); color: #0C4A6E; 
                border: 1px solid #7DD3FC; border-radius: 6px; 
                padding: 8px; font-size: 14px; font-family: 'Segoe UI', 'Microsoft YaHei', sans-serif; 
            }}
            QLineEdit:focus, QTextEdit:focus, QComboBox:focus {{ border: 2px solid #0284C7; background-color: #FFFFFF; }}
            
            QPushButton {{ color: white; border: none; border-radius: 6px; padding: 8px 12px; font-size: 14px; font-weight: bold; font-family: 'Segoe UI', 'Microsoft YaHei'; }}
            QPushButton:disabled {{ background-color: rgba(148, 163, 184, 0.8) !important; color: #F1F5F9; }}
            QPushButton#btn_primary {{ background-color: #0284C7; padding: 10px 15px; }}
            QPushButton#btn_primary:hover {{ background-color: #0369A1; }}
            QPushButton#btn_load {{ background-color: #0EA5E9; }}
            QPushButton#btn_load:hover {{ background-color: #0284C7; }}
            QPushButton.btn_general {{ background-color: #F0F9FF; color: #0284C7; border: 1px solid #7DD3FC; }}
            QPushButton.btn_general:hover {{ background-color: #E0F2FE; color: #0369A1; }}
            QPushButton.btn_math {{ background-color: #E0F2FE; color: #0284C7; border: 1px solid #38BDF8; }}
            QPushButton.btn_math:hover {{ background-color: #BAE6FD; }}
            QPushButton.btn_human {{ background-color: #FFF1F2; color: #E11D48; border: 1px solid #FECDD3; }}
            QPushButton.btn_human:hover {{ background-color: #FFE4E6; }}
            QPushButton.btn_english {{ background-color: #F0FDF4; color: #059669; border: 1px solid #BBF7D0; }}
            QPushButton.btn_english:hover {{ background-color: #DCFCE7; }}
            QLabel {{ font-size: 14px; font-weight: bold; color: #0C4A6E; background: transparent; }}
        """)

    def init_ui(self):
        main_widget = QWidget()
        self.setCentralWidget(main_widget)
        layout = QVBoxLayout(main_widget)
        layout.setContentsMargins(15, 15, 15, 15)
        layout.setSpacing(10)

        # ================= 第一层 =================
        top_frame = QFrame()
        top_frame.setProperty("class", "PanelFrame")
        top_layout = QHBoxLayout(top_frame)
        top_layout.setContentsMargins(15, 10, 15, 10)
        
        self.key_input = QLineEdit()
        self.key_input.setPlaceholderText("🔑 填入 API Key")
        self.key_input.setEchoMode(QLineEdit.EchoMode.Password)
        top_layout.addWidget(QLabel("API Key:"))
        top_layout.addWidget(self.key_input)
        
        top_layout.addWidget(QLabel(" 魔力引擎:"))
        self.model_combo = QComboBox()
        self.model_combo.setEditable(True) 
        self.model_combo.setMinimumWidth(180)
        default_models = [
            "deepseek-v3.2-thinking",
            "deepseek-v3.2",
            "gpt-4o-mini",
            "gpt-4o",
            "gpt-5.4",
            "claude-3-5-sonnet-20241022",
            "gemini-2.5-pro",
            "qwen3.5-plus",
            "kimi-k2-thinking"
        ]
        self.model_combo.addItems(default_models)
        top_layout.addWidget(self.model_combo)

        self.btn_load_ppt = QPushButton("🧊 冻结并导入课件")
        self.btn_load_ppt.setObjectName("btn_load")
        self.btn_load_ppt.clicked.connect(self.load_ppt)
        top_layout.addWidget(self.btn_load_ppt)
        layout.addWidget(top_frame)

        # ================= 第二层 =================
        gen_frame = QFrame()
        gen_frame.setProperty("class", "PanelFrame")
        gen_layout = QHBoxLayout(gen_frame)
        gen_layout.setContentsMargins(15, 10, 15, 10)
        gen_layout.addWidget(QLabel("🏫 完美教室:"))
        self.btn_summarize = QPushButton("✨ 冰雪总结")
        self.btn_summarize.setProperty("class", "btn_general")
        self.btn_summarize.clicked.connect(self.summarize_ppt)
        gen_layout.addWidget(self.btn_summarize)
        self.btn_concept = QPushButton("🧠 ⑨的核心闪卡")
        self.btn_concept.setProperty("class", "btn_general")
        self.btn_concept.clicked.connect(self.extract_concepts)
        gen_layout.addWidget(self.btn_concept)
        self.btn_quiz = QPushButton("📝 算数小测验")
        self.btn_quiz.setProperty("class", "btn_general")
        self.btn_quiz.clicked.connect(self.generate_quiz)
        gen_layout.addWidget(self.btn_quiz)
        self.btn_export = QPushButton("💾 导出天才笔记")
        self.btn_export.setProperty("class", "btn_general")
        self.btn_export.clicked.connect(self.export_notes)
        gen_layout.addWidget(self.btn_export)
        gen_layout.addStretch()
        gen_layout.addWidget(QLabel("投影:"))
        self.btn_zoom_out = QPushButton("🔍 缩小")
        self.btn_zoom_out.setProperty("class", "btn_general")
        self.btn_zoom_out.clicked.connect(lambda: self.adjust_zoom(-0.2))
        gen_layout.addWidget(self.btn_zoom_out)
        self.btn_zoom_in = QPushButton("🔍 放大")
        self.btn_zoom_in.setProperty("class", "btn_general")
        self.btn_zoom_in.clicked.connect(lambda: self.adjust_zoom(0.2))
        gen_layout.addWidget(self.btn_zoom_in)
        layout.addWidget(gen_frame)

        # ================= 第三层 =================
        subj_frame = QFrame()
        subj_frame.setProperty("class", "PanelFrame")
        subj_layout = QHBoxLayout(subj_frame)
        subj_layout.setContentsMargins(15, 10, 15, 10)
        subj_layout.addWidget(QLabel("📚 最强奥义:"))
        self.btn_variation = QPushButton("📐 算数: 绝对冰冻变式题")
        self.btn_variation.setProperty("class", "btn_math")
        self.btn_variation.clicked.connect(self.generate_variation)
        subj_layout.addWidget(self.btn_variation)
        self.btn_mindmap = QPushButton("🗺️ 文科: 幻想乡导图")
        self.btn_mindmap.setProperty("class", "btn_human")
        self.btn_mindmap.clicked.connect(self.generate_mindmap)
        subj_layout.addWidget(self.btn_mindmap)
        self.btn_vocab = QPushButton("🔠 英语: 妖精的词汇卡")
        self.btn_vocab.setProperty("class", "btn_english")
        self.btn_vocab.clicked.connect(self.generate_vocab)
        subj_layout.addWidget(self.btn_vocab)
        subj_layout.addStretch()
        layout.addWidget(subj_frame)

        # ================= 显示区 (搭载极速渲染引擎) =================
        self.chat_display = QWebEngineView()
        self.chat_display.setStyleSheet("background: transparent; border: none;")
        self.init_browser_engine()
        layout.addWidget(self.chat_display, stretch=1)

        # ================= 底部输入 =================
        bottom_frame = QFrame()
        bottom_frame.setProperty("class", "PanelFrame")
        bottom_layout = QHBoxLayout(bottom_frame)
        bottom_layout.setContentsMargins(10, 10, 10, 10)
        self.question_input = QTextEdit()
        self.question_input.setPlaceholderText("人类，有什么不懂的尽管问本天才！(Shift+Enter换行)...")
        self.question_input.setMaximumHeight(70)
        bottom_layout.addWidget(self.question_input)
        self.btn_ask = QPushButton("发射冰符 🚀")
        self.btn_ask.setObjectName("btn_primary")
        self.btn_ask.setMinimumHeight(70)
        self.btn_ask.setMinimumWidth(120)
        self.btn_ask.clicked.connect(self.ask_question)
        bottom_layout.addWidget(self.btn_ask)
        layout.addWidget(bottom_frame)

        self.control_buttons(False)

    # ❄️ 极速前端底座：加载一次，永不刷新
    def init_browser_engine(self):
        html_base = """
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="utf-8">
            <script src="https://cdn.staticfile.net/marked/12.0.1/marked.min.js"></script>
            <script>
                window.MathJax = {
                    tex: { inlineMath: [['$', '$'], ['\\(', '\\)']], displayMath: [['$$', '$$'], ['\\[', '\\]']] },
                    svg: { fontCache: 'global' }
                };
            </script>
            <script type="text/javascript" id="MathJax-script" async src="https://cdn.staticfile.net/mathjax/3.2.2/es5/tex-svg.min.js"></script>
            <script src="https://cdn.staticfile.net/mermaid/9.4.3/mermaid.min.js"></script>
            <script> mermaid.initialize({ startOnLoad: false, theme: 'default' }); </script>
            <style>
                body { font-family: 'Segoe UI', 'Microsoft YaHei', sans-serif; background: transparent; padding: 20px; margin: 0; color: #0F172A; overflow-x: hidden; }
                #chat-container { display: flex; flex-direction: column; gap: 24px; padding-bottom: 30px; }
                .message-row { display: flex; width: 100%; animation: fadeIn 0.3s ease-out forwards; }
                @keyframes fadeIn { from { opacity: 0; transform: translateY(10px); } to { opacity: 1; transform: translateY(0); } }
                .message-row.user { justify-content: flex-end; }
                .message-row.assistant, .message-row.system { justify-content: flex-start; }
                .avatar { width: 45px; height: 45px; border-radius: 50%; display: flex; align-items: center; justify-content: center; margin: 0 12px; flex-shrink: 0; box-shadow: 0 2px 5px rgba(0,0,0,0.1); font-size: 14px; font-weight: bold; color: white; text-align: center; }
                .user .avatar { background-color: #64748B; order: 2; }
                .assistant .avatar, .system .avatar { background-color: #0284C7; order: 1; border: 2px solid #BAE6FD; }
                .bubble { max-width: 85%; padding: 18px 25px; border-radius: 12px; font-size: 16px; line-height: 1.8; box-shadow: 0 4px 12px rgba(0,0,0,0.08); word-wrap: break-word; overflow-x: auto; }
                .user .bubble { background-color: rgba(241, 245, 249, 0.95); order: 1; color: #334155; border: 1px solid #CBD5E1; border-top-right-radius: 2px; }
                .assistant .bubble, .system .bubble { background-color: rgba(255, 255, 255, 0.95); border: 2px solid #BAE6FD; order: 2; border-top-left-radius: 2px; }
                h1, h2, h3 { color: #0284C7; border-bottom: 2px dashed #BAE6FD; padding-bottom: 6px; margin-top: 24px; }
                h1:first-child, h2:first-child, h3:first-child { margin-top: 0; }
                .bubble p { margin: 0 0 12px 0; }
                .bubble p:last-child { margin: 0; }
                .bubble code:not(.language-mermaid) { background-color: #F0F9FF; color: #0284C7; padding: 2px 6px; border-radius: 4px; font-family: Consolas, monospace; font-size: 0.95em; border: 1px solid #BAE6FD; }
                .bubble pre:not(:has(.language-mermaid)) { background-color: #0F172A; color: #E0F2FE; padding: 16px; border-radius: 8px; overflow-x: auto; font-size: 15px; border-left: 5px solid #38BDF8; }
                .bubble table { border-collapse: collapse; width: 100%; margin: 15px 0; font-size: 15px; }
                .bubble th, .bubble td { border: 1px solid #BAE6FD; padding: 12px; text-align: left; }
                .bubble th { background-color: #F0F9FF; font-weight: bold; color: #0284C7; }
                .bubble blockquote { border-left: 4px solid #38BDF8; background-color: #F0F9FF; margin: 15px 0; padding: 12px 18px; color: #0369A1; font-style: italic; border-radius: 0 8px 8px 0; }
                .mermaid-container { background-color: #FFFFFF; border: 2px dashed #38BDF8; border-radius: 8px; padding: 20px; text-align: center; margin: 15px 0; }
                .MathJax { font-size: 110% !important; }
                
                /* 加载中动画特效 */
                .loading-text { font-style: italic; color: #0284C7; animation: blink 1.5s infinite; }
                @keyframes blink { 0% { opacity: 0.4; } 50% { opacity: 1; } 100% { opacity: 0.4; } }
            </style>
        </head>
        <body>
            <div id="chat-container"></div>
            <script>
                let msgCounter = 0;
                
                // 页面加载完自动问好
                window.onload = () => {
                    const welcome = "### ❄️ 欢迎来到《琪露诺的完美算数教室》！\\n\\n本天才⑨（最强教研AI）将亲自为你解答！\\n\\n1. 填入你的 API Key，导入课件！\\n2. 请在上方自由选择想使用的 AI 模型。\\n3. **【已搭载极速增量引擎】**，点击按钮**瞬间响应**，绝不卡顿！";
                    appendMessage('system', welcome);
                };

                // 清空黑板
                function clearChat() {
                    document.getElementById('chat-container').innerHTML = '';
                    msgCounter = 0;
                }

                // 追加新气泡（瞬间完成，不重绘旧元素）
                function appendMessage(role, rawContent) {
                    const loadingEl = document.getElementById('loading-row');
                    if(loadingEl) loadingEl.remove();

                    const container = document.getElementById('chat-container');
                    const row = document.createElement('div');
                    row.className = `message-row ${role}`;
                    
                    const avatar = document.createElement('div');
                    avatar.className = 'avatar';
                    avatar.innerHTML = role === 'user' ? '人类' : '天才⑨';
                    
                    const bubble = document.createElement('div');
                    bubble.className = 'bubble';
                    bubble.innerHTML = marked.parse(rawContent);
                    
                    const mermaidBlocks = bubble.querySelectorAll('code.language-mermaid');
                    mermaidBlocks.forEach((block, idx) => {
                        const graphCode = block.textContent;
                        const pre = block.parentElement;
                        const div = document.createElement('div');
                        div.className = 'mermaid-container';
                        div.id = `mermaid-${msgCounter}-${idx}`; 
                        pre.replaceWith(div);
                        try {
                            mermaid.mermaidAPI.render(`svg-${msgCounter}-${idx}`, graphCode, (svgCode) => {
                                div.innerHTML = svgCode;
                            });
                        } catch (e) {
                            div.innerHTML = `<p style="color:red;">哎呀，魔法回路出错了！</p><pre>${graphCode}</pre>`;
                        }
                    });
                    
                    row.appendChild(avatar);
                    row.appendChild(bubble);
                    container.appendChild(row);
                    msgCounter++;

                    if (window.MathJax && window.MathJax.typesetPromise) {
                        MathJax.typesetPromise([bubble]).then(() => {
                            window.scrollTo({top: document.body.scrollHeight, behavior: 'smooth'});
                        });
                    } else {
                        window.scrollTo({top: document.body.scrollHeight, behavior: 'smooth'});
                    }
                }

                // 瞬间弹出等待动画
                function showLoading(modelName) {
                    const container = document.getElementById('chat-container');
                    const row = document.createElement('div');
                    row.className = `message-row assistant`;
                    row.id = 'loading-row';
                    
                    const avatar = document.createElement('div');
                    avatar.className = 'avatar';
                    avatar.innerHTML = '天才⑨';
                    
                    const bubble = document.createElement('div');
                    bubble.className = 'bubble loading-text';
                    bubble.innerHTML = `❄️ 天才的大脑 [${modelName}] 正在极速推演中，请稍候...`;
                    
                    row.appendChild(avatar);
                    row.appendChild(bubble);
                    container.appendChild(row);
                    window.scrollTo({top: document.body.scrollHeight, behavior: 'smooth'});
                }
            </script>
        </body>
        </html>
        """
        self.chat_display.setHtml(html_base)

    # ================= 网页交互接口 =================
    def js_append(self, role, content):
        safe_content = json.dumps(content)
        self.chat_display.page().runJavaScript(f"appendMessage('{role}', {safe_content});")

    def js_show_loading(self, model_name):
        safe_model = json.dumps(model_name)
        self.chat_display.page().runJavaScript(f"showLoading({safe_model});")
        
    def js_clear_chat(self):
        self.chat_display.page().runJavaScript("clearChat();")

    def control_buttons(self, state):
        self.btn_summarize.setEnabled(state)
        self.btn_concept.setEnabled(state)
        self.btn_quiz.setEnabled(state)
        self.btn_export.setEnabled(state)
        self.btn_zoom_in.setEnabled(state)
        self.btn_zoom_out.setEnabled(state)
        self.btn_variation.setEnabled(state)
        self.btn_mindmap.setEnabled(state)
        self.btn_vocab.setEnabled(state)
        self.btn_ask.setEnabled(state)

    def adjust_zoom(self, delta):
        self.zoom_factor += delta
        self.zoom_factor = max(0.5, min(self.zoom_factor, 3.0))
        self.chat_display.setZoomFactor(self.zoom_factor)

    def load_ppt(self):
        file_path, _ = QFileDialog.getOpenFileName(self, "选择课件", "", "PowerPoint Files (*.pptx)")
        if not file_path:
            return
        try:
            prs = Presentation(file_path)
            text_runs = []
            for slide_idx, slide in enumerate(prs.slides):
                text_runs.append(f"\n[第 {slide_idx + 1} 页]")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_runs.append(shape.text.strip())
            
            self.ppt_content = "\n".join(text_runs)
            if len(self.ppt_content) > 10000:
                self.ppt_content = self.ppt_content[:10000] + "\n...(内容过长已截断)"

            QMessageBox.information(self, "冰冻完成", f"❄️ 成功提取课件文字！\n现在，所有的魔法大招都为你敞开！")
            self.control_buttons(True)
        except Exception as e:
            QMessageBox.critical(self, "错误", f"解析失败: {str(e)}")

    # ================= 核心多线程调度系统 =================
    def call_ai_async(self, user_prompt, is_summary=False):
        self.api_key = self.key_input.text().strip()
        if not self.api_key:
            QMessageBox.warning(self, "警告", "没有魔力来源 (API Key) 是不行的！")
            return

        selected_model = self.model_combo.currentText().strip()
        if not selected_model:
            selected_model = "deepseek-v3.2-thinking"
            self.model_combo.setCurrentText(selected_model)

        base_url = "https://api.chatanywhere.com.cn/v1"

        system_prompt = (
            "你是《琪露诺的完美算数教室》的首席教研导师。你极其聪明、专业且严谨，但在回答的开头或结尾可以稍微带有一点骄傲可爱的语气（因为你自称天才⑨）。\n"
            "无论如何，知识点的讲解必须达到全国重点高中的最高学术标准，绝不含糊。\n"
            "数学公式必须用 LaTeX（行内 $...$，独立 $$...$$）。结构图必须提供 Mermaid 代码块。\n"
            f"这是今天的课件：\n{self.ppt_content}"
        )

        if is_summary:
            self.chat_history = []
            self.js_clear_chat()
            self.js_append("system", "### ❄️ 全新推演开始！\n\n本天才⑨已清空之前的黑板，准备好接受知识的洗礼吧！")
        
        messages = [{"role": "system", "content": system_prompt}]
        messages.extend(self.chat_history)
        messages.append({"role": "user", "content": user_prompt})

        # 禁用按钮，瞬间渲染用户消息和加载动画
        self.control_buttons(False)
        self.js_append("user", user_prompt)
        self.js_show_loading(selected_model)

        self.current_user_prompt = user_prompt

        # 启动后台引擎
        self.worker = AIWorker(self.api_key, base_url, selected_model, messages)
        self.worker.finished_signal.connect(self.on_ai_success)
        self.worker.error_signal.connect(self.on_ai_error)
        self.worker.start()

    def on_ai_success(self, ai_reply):
        self.chat_history.append({"role": "user", "content": self.current_user_prompt})
        self.chat_history.append({"role": "assistant", "content": ai_reply})
        
        self.js_append("assistant", ai_reply)
        self.control_buttons(True)

    def on_ai_error(self, error_msg):
        self.js_append("system", f"❌ **发生网络错误或 API 受限**：\n```text\n{error_msg}\n```")
        self.control_buttons(True)

    # ================= 各项功能调用 =================
    def summarize_ppt(self):
        prompt = "请对这份PPT进行【全局精简总结】。要求：\n1. 使用 Markdown 标题（#，##）分层级列出核心框架。\n2. 语言极度精炼，只保留核心定理、公式和关键脉络，去掉废话。\n3. 如果有公式，务必用 LaTeX 语法打出。"
        self.call_ai_async(prompt, is_summary=True)

    def extract_concepts(self):
        prompt = "提取出 5 个最核心的考点。请用一点点骄傲但又确实解答得很完美的口吻，生成通俗易懂的“课堂闪卡”解释。"
        self.call_ai_async(prompt)

    def generate_quiz(self):
        prompt = "根据PPT内容，出 3 道单选题作为随堂测验。每道题必须有 A/B/C/D 四个选项。最后统一给出【正确答案】和【详细解析】。"
        self.call_ai_async(prompt)

    def generate_variation(self):
        prompt = "提取核心理科例题。设计一道达到【全国卷高考难度】的原创压轴变式题。提供详细的阶梯式参考答案。公式严格使用 LaTeX。"
        self.call_ai_async(prompt)

    def generate_mindmap(self):
        prompt = """对这份PPT进行结构化解构，提取核心脉络。严格使用 Mermaid 语法 (mindmap) 绘制思维导图代码。代码必须放在 ```mermaid 的代码块中。"""
        self.call_ai_async(prompt)

    def generate_vocab(self):
        prompt = """挖掘PPT底层知识点，提取 5 个【高中必备英语核心词汇】。用排版精美的 Markdown 列表输出：英文、词性、音标、中文、英文例句及翻译。"""
        self.call_ai_async(prompt)

    def ask_question(self):
        question = self.question_input.toPlainText().strip()
        if not question:
            return
        self.question_input.clear()
        self.call_ai_async(question)

    def export_notes(self):
        if not self.chat_history:
            QMessageBox.information(self, "提示", "黑板上还没写字呢，无法导出！")
            return
        file_path, _ = QFileDialog.getSaveFileName(self, "💾 导出天才笔记", "琪露诺的完美笔记.md", "Markdown Files (*.md)")
        if file_path:
            try:
                with open(file_path, 'w', encoding='utf-8') as f:
                    f.write("# ❄️ 琪露诺的完美算数教室 - 专属笔记\n\n> 哪怕是⑨也能看懂的完美解析！\n\n---\n\n")
                    for msg in self.chat_history:
                        if msg['role'] == 'user':
                            f.write(f"### 🧑 人类提问：\n{msg['content']}\n\n")
                        elif msg['role'] == 'assistant':
                            f.write(f"### ❄️ 天才⑨解答：\n{msg['content']}\n\n---\n\n")
                QMessageBox.information(self, "大成功", "笔记导出成功！快拿去给学生们复习吧！")
            except Exception as e:
                pass

if __name__ == "__main__":
    app = QApplication(sys.argv)
    window = PPTAIAssistant()
    window.show()
    sys.exit(app.exec())
